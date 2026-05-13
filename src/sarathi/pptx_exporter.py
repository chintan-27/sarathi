from __future__ import annotations

import re
from pathlib import Path

# ── Public API ────────────────────────────────────────────────────────────────

def to_pptx(outline: dict, files_map: dict, output_path: Path) -> None:
    """Convert a builder JSON outline to a .pptx file using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    for slide_def in outline.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)
        _apply_background(slide, outline.get("theme", "dark-gradient"))
        _render_slide(slide, slide_def, files_map, prs)

    prs.save(str(output_path))


# ── Theme colours ─────────────────────────────────────────────────────────────

_THEME_COLORS = {
    "dark-gradient": {
        "bg":     (13,  13,  26),
        "fg":     (232, 232, 240),
        "accent": (79,  195, 247),
        "dim":    (136, 136, 170),
    },
    "dracula": {
        "bg":     (40,  42,  54),
        "fg":     (248, 248, 242),
        "accent": (189, 147, 249),
        "dim":    (98,  114, 164),
    },
    "light": {
        "bg":     (250, 250, 250),
        "fg":     (26,  26,  46),
        "accent": (21,  101, 192),
        "dim":    (102, 102, 102),
    },
    "minimal": {
        "bg":     (17,  17,  17),
        "fg":     (238, 238, 238),
        "accent": (255, 255, 255),
        "dim":    (136, 136, 136),
    },
}


def _colors(theme: str) -> dict:
    return _THEME_COLORS.get(theme, _THEME_COLORS["dark-gradient"])


def _rgb(triplet: tuple) -> "RGBColor":
    from pptx.dml.color import RGBColor
    return RGBColor(*triplet)


# ── Background fill ───────────────────────────────────────────────────────────

def _apply_background(slide, theme: str) -> None:
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    c = _colors(theme)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(c["bg"])


# ── Slide rendering ───────────────────────────────────────────────────────────

def _render_slide(slide, slide_def: dict, files_map: dict, prs) -> None:
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    theme = prs.slide_width  # we pass theme separately via outline
    # Re-derive theme from the prs object isn't straightforward; use default
    # The outline dict carries theme info at the top level — accessed via slide_def
    theme_name = slide_def.get("_theme", "dark-gradient")
    c = _colors(theme_name)

    stype    = slide_def.get("type", "context")
    heading  = slide_def.get("heading", "")
    insight  = slide_def.get("insight", "")
    artifacts = slide_def.get("artifacts", [])
    hero     = slide_def.get("hero_metric") or slide_def.get("layout_hint", "")
    layout   = slide_def.get("layout_hint", "")

    W = Inches(13.33)
    H = Inches(7.5)

    if stype == "title":
        _title_slide(slide, heading, insight, c, W, H)
    elif stype == "metric_callout":
        _metric_slide(slide, heading, hero or insight, c, W, H)
    elif stype in ("chart", "image"):
        _image_slide(slide, heading, artifacts, files_map, c, W, H)
    elif stype == "code":
        _code_slide(slide, heading, artifacts, files_map, c, W, H)
    elif stype == "takeaways":
        _bullets_slide(slide, heading, insight, c, W, H, accent=True)
    else:
        _content_slide(slide, heading, insight, c, W, H)

    # Speaker notes
    notes = slide_def.get("speaker_notes", "")
    if notes:
        tf = slide.notes_slide.notes_text_frame
        tf.text = notes


# ── Slide type renderers ──────────────────────────────────────────────────────

def _add_textbox(slide, text: str, left, top, width, height,
                 font_size: int, bold: bool, color: tuple,
                 align="left") -> None:
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    from pptx.dml.color import RGBColor
    run.font.color.rgb = RGBColor(*color)


def _title_slide(slide, heading, subtitle, c, W, H) -> None:
    from pptx.util import Inches, Pt
    pad = Inches(0.8)
    _add_textbox(slide, heading,
                 pad, H * 0.3, W - pad * 2, Inches(1.5),
                 font_size=44, bold=True, color=c["accent"], align="center")
    if subtitle:
        _add_textbox(slide, subtitle,
                     pad, H * 0.55, W - pad * 2, Inches(1.2),
                     font_size=20, bold=False, color=c["dim"], align="center")


def _metric_slide(slide, heading, metric_text, c, W, H) -> None:
    from pptx.util import Inches, Pt
    pad = Inches(0.8)
    _add_textbox(slide, heading,
                 pad, Inches(0.4), W - pad * 2, Inches(0.8),
                 font_size=24, bold=False, color=c["dim"], align="center")
    # Big number centred
    clean = re.sub(r"<[^>]+>", "", metric_text).strip()[:80]
    _add_textbox(slide, clean,
                 pad, H * 0.25, W - pad * 2, H * 0.5,
                 font_size=80, bold=True, color=c["accent"], align="center")


def _image_slide(slide, heading, artifacts, files_map, c, W, H) -> None:
    from pptx.util import Inches, Pt
    import io, base64
    pad = Inches(0.5)
    _add_textbox(slide, heading,
                 pad, Inches(0.15), W - pad * 2, Inches(0.65),
                 font_size=26, bold=True, color=c["accent"])

    for art_path in artifacts:
        rf = files_map.get(art_path) or files_map.get(Path(art_path).name)
        if rf and rf.type in ("image", "svg") and rf.content.startswith("data:"):
            try:
                _, b64 = rf.content.split(",", 1)
                img_bytes = base64.b64decode(b64)
                img_io = io.BytesIO(img_bytes)
                slide.shapes.add_picture(
                    img_io,
                    left=pad,
                    top=Inches(0.9),
                    width=W - pad * 2,
                    height=H - Inches(1.1),
                )
                break
            except Exception:
                pass

    if not any(rf and rf.type in ("image", "svg")
               for art_path in artifacts
               for rf in [files_map.get(art_path) or files_map.get(Path(art_path).name)]
               if rf):
        _add_textbox(slide, heading,
                     pad, H * 0.3, W - pad * 2, H * 0.4,
                     font_size=20, bold=False, color=c["fg"], align="center")


def _code_slide(slide, heading, artifacts, files_map, c, W, H) -> None:
    from pptx.util import Inches, Pt
    pad = Inches(0.5)
    _add_textbox(slide, heading,
                 pad, Inches(0.15), W - pad * 2, Inches(0.65),
                 font_size=26, bold=True, color=c["accent"])

    code_text = ""
    for art_path in artifacts:
        rf = files_map.get(art_path) or files_map.get(Path(art_path).name)
        if rf and rf.type == "code":
            code_text = rf.content[:1200]
            break

    if code_text:
        from pptx.util import Pt
        txBox = slide.shapes.add_textbox(
            pad, Inches(0.9), W - pad * 2, H - Inches(1.1)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = code_text
        run = p.runs[0]
        run.font.size = Pt(13)
        from pptx.dml.color import RGBColor
        run.font.color.rgb = RGBColor(*c["fg"])
        run.font.name = "Courier New"


def _bullets_slide(slide, heading, insight, c, W, H, accent=False) -> None:
    from pptx.util import Inches, Pt
    pad = Inches(0.6)
    _add_textbox(slide, heading,
                 pad, Inches(0.2), W - pad * 2, Inches(0.75),
                 font_size=30, bold=True,
                 color=c["accent"] if accent else c["fg"])

    bullets = [s.strip() for s in re.split(r"[•\-\n]", insight) if s.strip()][:6]
    if not bullets and insight:
        bullets = [insight]

    txBox = slide.shapes.add_textbox(pad, Inches(1.1), W - pad * 2, H - Inches(1.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    from pptx.dml.color import RGBColor

    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"• {bullet}"
        from pptx.util import Pt
        if p.runs:
            p.runs[0].font.size = Pt(20)
            p.runs[0].font.color.rgb = RGBColor(*c["fg"])
        p.space_after = Pt(8)


def _content_slide(slide, heading, insight, c, W, H) -> None:
    from pptx.util import Inches, Pt
    pad = Inches(0.6)
    _add_textbox(slide, heading,
                 pad, Inches(0.2), W - pad * 2, Inches(0.8),
                 font_size=30, bold=True, color=c["accent"])
    if insight:
        _add_textbox(slide, insight,
                     pad, Inches(1.15), W - pad * 2, H - Inches(1.4),
                     font_size=20, bold=False, color=c["fg"])
